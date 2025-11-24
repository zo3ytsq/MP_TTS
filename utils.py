from pptx import Presentation

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    slides_text = []

    for slide in prs.slides:
        text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                line = shape.text.strip()
                # Skip common footer text or copyright
                if (
                    not line
                    or "©" in line
                    or "copyright" in line.lower()
                    or "temasek polytechnic" in line.lower()
                ):
                    continue
                text.append(line)
        slides_text.append(text)

    return slides_text